"""Build a 3-slide SPI/EWS set in the dark-navy Canva style (UTY crest top-left,
purple label block + big white title on the left, white rounded content card on
the right). Output: spi-canva.pptx -- importable into Canva.

Run:  uv run --with python-pptx --with pillow python build_spi_canva.py
Content source: panduan-demo-sidang.html (sections RAG / EWS / code).
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
LOGO = os.path.join(HERE, "assets", "uty-logo.png")

# ---- palette (dark Canva) ----
BG     = RGBColor(0x0A, 0x0E, 0x20)   # deep navy-black
GLOW   = RGBColor(0x1B, 0x36, 0x7A)   # blue glow
BAR    = RGBColor(0x2D, 0x55, 0xC8)   # right accent bar
INDIGO = RGBColor(0x49, 0x3A, 0xA8)   # purple label block
TEAL   = RGBColor(0x35, 0x9E, 0x8F)   # chevron
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CARD   = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x12, 0x22, 0x40)   # card heading
BODY   = RGBColor(0x3A, 0x47, 0x5E)   # card body
MUTED  = RGBColor(0x8A, 0x96, 0xAC)
BLUE   = RGBColor(0x24, 0x5E, 0xD6)
TINTB  = RGBColor(0xEC, 0xF1, 0xFD)   # blue tint
LINE   = RGBColor(0xDD, 0xE5, 0xF1)
GREEN  = RGBColor(0x1E, 0x9E, 0x57); GTINT = RGBColor(0xE7, 0xF5, 0xEC)
AMBER  = RGBColor(0xC6, 0x86, 0x00); ATINT = RGBColor(0xFB, 0xF1, 0xDB)
RED    = RGBColor(0xD3, 0x3A, 0x2C); RTINT = RGBColor(0xFB, 0xE9, 0xE7)

TITLEF = "Arial Black"
HEADF  = "Arial"
BODYF  = "Calibri"
MONOF  = "Consolas"


def make_logo():
    """Key out the near-black background of the naskah cover crest -> transparent PNG."""
    if os.path.exists(LOGO):
        return
    from PIL import Image
    src = os.path.join(HERE, "..", "naskah-extract", "images", "page-001-img00.jpeg")
    im = Image.open(src).convert("RGBA")
    px = im.load()
    w, h = im.size
    for y in range(h):
        for x in range(w):
            r, g, b, _ = px[x, y]
            if max(r, g, b) < 40:           # deep black surround -> transparent
                px[x, y] = (r, g, b, 0)
    # tight crop to content
    bbox = im.getbbox()
    if bbox:
        im = im.crop(bbox)
    os.makedirs(os.path.dirname(LOGO), exist_ok=True)
    im.save(LOGO)


def set_alpha(shape, pct):
    """Set fill transparency: pct = opacity percent (0..100)."""
    sp = shape.fill._xPr.find(qn("a:solidFill"))
    clr = sp.find(qn("a:srgbClr"))
    a = clr.makeelement(qn("a:alpha"), {"val": str(int(pct * 1000))})
    clr.append(a)


def no_line(shape):
    shape.line.fill.background()


def rect(slide, x, y, w, h, fill, line=None, lw=1.0, rounded=False, rad=0.12, shadow=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        try: shp.adjustments[0] = rad
        except Exception: pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        no_line(shp)
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def oval(slide, x, y, w, h, fill, line=None, lw=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: no_line(shp)
    else: shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def chevron(slide, x, y, s, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(s), Inches(s))
    shp.fill.solid(); shp.fill.fore_color.rgb = color; no_line(shp); shp.shadow.inherit = False
    return shp


def txt(slide, x, y, w, h, runs, size, color=INK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font=BODYF, spacing=None, line_spacing=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    if isinstance(runs, str):
        runs = [(runs, color, bold)]
    p = tf.paragraphs[0]; p.alignment = align
    if line_spacing: p.line_spacing = line_spacing
    for t, c, b in runs:
        r = p.add_run(); r.text = t
        r.font.size = Pt(size); r.font.bold = b; r.font.name = font
        r.font.color.rgb = c
        if spacing is not None:
            rPr = r._r.get_or_add_rPr(); rPr.set("spc", str(int(spacing)))
    return tb


def para(tb, runs, size, color=BODY, bold=False, align=PP_ALIGN.LEFT, font=BODYF, space_before=4):
    p = tb.text_frame.add_paragraph(); p.alignment = align; p.space_before = Pt(space_before)
    if isinstance(runs, str): runs = [(runs, color, bold)]
    for t, c, b in runs:
        r = p.add_run(); r.text = t; r.font.size = Pt(size); r.font.bold = b
        r.font.name = font; r.font.color.rgb = c
    return p


def base_slide(prs, label, title_lines, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # background
    rect(s, 0, 0, 13.333, 7.5, BG)
    g = oval(s, -3.2, -3.6, 8.5, 8.5, GLOW); set_alpha(g, 22)
    rect(s, 12.95, 0, 0.38, 7.5, BAR)
    gb = rect(s, 12.55, 0, 0.4, 7.5, BLUE); set_alpha(gb, 20)
    # header: crest + uni name
    if os.path.exists(LOGO):
        s.shapes.add_picture(LOGO, Inches(0.55), Inches(0.42), height=Inches(0.62))
    txt(s, 1.32, 0.46, 6.0, 0.6,
        [("UNIVERSITAS TEKNOLOGI YOGYAKARTA", WHITE, True)],
        13, anchor=MSO_ANCHOR.MIDDLE, font=HEADF, spacing=60)
    # left column: label block, big title, chevron, subtitle
    lblw = 0.55 + len(label) * 0.135
    rect(s, 0.6, 2.0, lblw, 0.5, INDIGO, rounded=True, rad=0.5)
    txt(s, 0.6, 2.0, lblw, 0.5, [(label, WHITE, True)], 14,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=HEADF, spacing=30)
    ty = 2.72
    tt = txt(s, 0.55, ty, 4.7, 2.0,
             [(title_lines[0], WHITE, True)], 42, font=TITLEF, line_spacing=0.95)
    for extra in title_lines[1:]:
        para(tt, [(extra, WHITE, True)], 42, font=TITLEF, space_before=0)
    chevron(s, 0.62, 4.95, 0.5, TEAL)
    chevron(s, 0.92, 4.95, 0.5, TEAL);
    txt(s, 0.6, 5.75, 4.4, 1.0, [(subtitle, RGBColor(0xB9,0xC6,0xE0), False)], 15,
        font=BODYF, line_spacing=1.05)
    # page footer (subtle)
    return s


def card(slide, x=5.35, y=1.35, w=7.4, h=4.9):
    bk = rect(slide, x + 0.07, y + 0.11, w, h, RGBColor(0x00, 0x00, 0x00), rounded=True, rad=0.045)
    set_alpha(bk, 26)
    return rect(slide, x, y, w, h, CARD, rounded=True, rad=0.045)


def slide_formula(prs):
    s = base_slide(prs, "INTI SISTEM", ["ALGORITMA", "SPI"],
                   "Schedule Performance Index — satu angka objektif yang menggerakkan seluruh Early Warning System.")
    cx, cy, cw = 5.35, 1.35, 7.4
    card(s, cx, cy, cw, 4.9)
    pad = 0.45
    txt(s, cx + pad, cy + 0.32, cw - 2 * pad, 0.4, [("RUMUS INTI", BLUE, True)], 13, font=HEADF, spacing=40)
    # big formula
    fb = rect(s, cx + pad, cy + 0.78, cw - 2 * pad, 1.05, TINTB, rounded=True, rad=0.12)
    txt(s, cx + pad, cy + 0.78, cw - 2 * pad, 1.05,
        [("SPI  =  EV  /  PV", INK, True)], 34, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE, font=TITLEF)
    # definitions
    dy = cy + 2.05
    txt(s, cx + pad, dy, cw - 2 * pad, 0.55,
        [("EV", GREEN, True), ("  Earned Value  —  % pekerjaan yang ", BODY, False),
         ("benar-benar selesai", INK, True)], 15, font=BODYF)
    txt(s, cx + pad, dy + 0.55, cw - 2 * pad, 0.7,
        [("PV", BLUE, True), ("  Planned Value  —  % yang ", BODY, False),
         ("seharusnya", INK, True), ("  selesai = (hari berjalan / total durasi) × 100%", BODY, False)],
        15, font=BODYF)
    rect(s, cx + pad, dy + 1.32, cw - 2 * pad, 0.012, LINE)
    # interpretation pills
    py = dy + 1.5
    pills = [("SPI = 1  ·  tepat jadwal", GTINT, GREEN),
             ("SPI > 1  ·  lebih cepat", TINTB, BLUE),
             ("SPI < 1  ·  terlambat", RTINT, RED)]
    px = cx + pad; pw = (cw - 2 * pad - 0.4) / 3
    for t, fc, tc in pills:
        rect(s, px, py, pw, 0.62, fc, rounded=True, rad=0.5)
        txt(s, px, py, pw, 0.62, [(t, tc, True)], 13, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE, font=BODYF)
        px += pw + 0.2


def rag_row(slide, x, y, w, name, cond, mean, tint, color):
    rect(slide, x, y, w, 0.92, tint, rounded=True, rad=0.14)
    oval(slide, x + 0.26, y + 0.31, 0.3, 0.3, color)
    txt(slide, x + 0.72, y, 1.7, 0.92, [(name, color, True)], 19,
        anchor=MSO_ANCHOR.MIDDLE, font=TITLEF)
    txt(slide, x + 2.5, y + 0.13, 2.0, 0.7,
        [(cond, INK, True)], 15, anchor=MSO_ANCHOR.MIDDLE, font=MONOF)
    txt(slide, x + 4.55, y, w - 4.8, 0.92, [(mean, BODY, False)], 14,
        anchor=MSO_ANCHOR.MIDDLE, font=BODYF)


def slide_rag(prs):
    s = base_slide(prs, "INDIKATOR", ["WARNA", "RAG"],
                   "Tiga warna kesehatan proyek. SPI diubah jadi warna supaya manajer paham dalam sekejap.")
    cx, cy, cw = 5.35, 1.35, 7.4
    card(s, cx, cy, cw, 4.9)
    pad = 0.42
    txt(s, cx + pad, cy + 0.3, cw - 2 * pad, 0.4, [("AMBANG BATAS (RAG)", BLUE, True)], 13, font=HEADF, spacing=40)
    ry = cy + 0.82; rw = cw - 2 * pad
    rag_row(s, cx + pad, ry, rw, "HIJAU", "SPI ≥ 0,95", "Sesuai jadwal (on track)", GTINT, GREEN)
    rag_row(s, cx + pad, ry + 1.02, rw, "KUNING", "0,85 – 0,95", "Perlu perhatian (warning)", ATINT, AMBER)
    rag_row(s, cx + pad, ry + 2.04, rw, "MERAH", "SPI < 0,85", "Kritis — tertinggal jadwal", RTINT, RED)
    # example strip
    ey = ry + 3.2
    txt(s, cx + pad, ey, cw - 2 * pad, 0.34,
        [("CONTOH — Villa Kaliurang, 6 tugas (tandai “Selesai” satu per satu):", MUTED, True)],
        12, font=BODYF)
    chips = [("1 → 0,22", RTINT, RED), ("3 → 0,67", RTINT, RED),
             ("4 → 0,89", ATINT, AMBER), ("5 → 1,11", GTINT, GREEN), ("6 → 1,33", GTINT, GREEN)]
    px = cx + pad; cy2 = ey + 0.4
    for t, fc, tc in chips:
        cwid = 1.18
        rect(s, px, cy2, cwid, 0.5, fc, rounded=True, rad=0.4)
        txt(s, px, cy2, cwid, 0.5, [(t, tc, True)], 13, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE, font=MONOF)
        px += cwid + 0.16


def flow_chip(slide, x, y, w, num, label, sub):
    rect(slide, x, y, w, 0.92, TINTB, rounded=True, rad=0.16)
    oval(slide, x + 0.13, y + 0.13, 0.34, 0.34, BLUE)
    txt(slide, x + 0.13, y + 0.13, 0.34, 0.34, [(str(num), WHITE, True)], 13,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=HEADF)
    txt(slide, x + 0.56, y + 0.12, w - 0.66, 0.42, [(label, INK, True)], 12.5, font=BODYF,
        anchor=MSO_ANCHOR.MIDDLE)
    txt(slide, x + 0.56, y + 0.5, w - 0.66, 0.34, [(sub, MUTED, False)], 9.5, font=MONOF)


def slide_flow(prs):
    s = base_slide(prs, "OTOMATIS", ["ALUR", "EWS"],
                   "Setiap status tugas berubah, SPI dihitung ulang otomatis — tanpa input manual.")
    cx, cy, cw = 5.35, 1.2, 7.4
    card(s, cx, cy, cw, 5.2)
    pad = 0.4
    txt(s, cx + pad, cy + 0.26, cw - 2 * pad, 0.4, [("ALUR OTOMATIS (8 LANGKAH)", BLUE, True)], 13, font=HEADF, spacing=40)
    steps = [("Status tugas", "berubah"), ("Hitung ulang", "recalculateSPI()"),
             ("Hitung EV & PV", "selesai / total"), ("SPI = EV / PV", "rasio"),
             ("SPI jadi warna", "categorizeHealth()"), ("Simpan", "project_health"),
             ("Dashboard", "merah teratas"), ("Manajer", "bertindak → dini")]
    colw = (cw - 2 * pad - 3 * 0.22) / 4
    y0 = cy + 0.74
    for i, (lab, sub) in enumerate(steps):
        row, col = divmod(i, 4)
        x = cx + pad + col * (colw + 0.22)
        y = y0 + row * 1.12
        flow_chip(s, x, y, colw, i + 1, lab, sub)
        if col < 3:
            txt(s, x + colw - 0.02, y, 0.24, 0.92, [("›", BLUE, True)], 20,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=HEADF)
    # 3 layers of warning
    ly = y0 + 2.45
    rect(s, cx + pad, ly, cw - 2 * pad, 0.012, LINE)
    txt(s, cx + pad, ly + 0.12, cw - 2 * pad, 0.34, [("3 LAPIS PERINGATAN", INK, True)], 12.5, font=HEADF, spacing=30)
    layers = [("Warna RAG", "per proyek"), ("Urut kritis", "merah di atas"), ("Ringkasan", "jumlah R/A/G")]
    lw = (cw - 2 * pad - 2 * 0.2) / 3; lx = cx + pad
    for lab, sub in layers:
        rect(s, lx, ly + 0.52, lw, 0.66, RGBColor(0xF4, 0xF7, 0xFC), rounded=True, rad=0.14, line=LINE)
        txt(s, lx + 0.16, ly + 0.52, lw - 0.3, 0.66,
            [(lab + "  ", INK, True), (sub, MUTED, False)], 12, anchor=MSO_ANCHOR.MIDDLE, font=BODYF)
        lx += lw + 0.2


def main():
    make_logo()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_formula(prs)
    slide_rag(prs)
    slide_flow(prs)
    out = os.path.join(HERE, "spi-canva.pptx")
    prs.save(out)
    print("saved", out)


if __name__ == "__main__":
    main()
