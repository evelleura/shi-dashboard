# /// script
# requires-python = ">=3.10"
# dependencies = ["python-pptx>=1.0", "lxml"]
# ///
"""Build the light-themed, animated sidang TA deck (.pptx) for Dian Putri Iswandi.
Run: uv run build_pptx.py  ->  presentasi-sidang.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import os

# ---------- palette (LIGHT) ----------
WHITE = RGBColor(0xFF,0xFF,0xFF)
BG    = RGBColor(0xFF,0xFF,0xFF)
CARD  = RGBColor(0xF2,0xF6,0xFC)
CARD2 = RGBColor(0xE9,0xF0,0xFA)
INK   = RGBColor(0x16,0x23,0x3A)
NAVY  = RGBColor(0x13,0x29,0x4B)
MUTED = RGBColor(0x5E,0x6E,0x86)
BLUE  = RGBColor(0x1F,0x6F,0xE0)
BLUED = RGBColor(0x16,0x4F,0xA8)
GREEN = RGBColor(0x1E,0x9E,0x57)
AMBER = RGBColor(0xCF,0x8A,0x00)
RED   = RGBColor(0xD3,0x3A,0x2C)
LINE  = RGBColor(0xD3,0xDE,0xEC)
GOLD  = RGBColor(0xC8,0x86,0x00)

HEAD = "Georgia"
BODY = "Calibri"
MONO = "Consolas"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ---------- helpers ----------
def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s

def rect(s, x,y,w,h, fill=CARD, line=LINE, lw=1.0, rounded=True, rad=0.06, shadow=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    if rounded:
        try: shp.adjustments[0] = rad
        except Exception: pass
    return shp

def oval(s, x,y,w,h, fill=CARD, line=LINE, lw=1.0):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp

def chev(s, x,y,w,h, color=MUTED):
    shp = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x),Inches(y),Inches(w),Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp

def dot(s, x,y, d, color):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x),Inches(y),Inches(d),Inches(d))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp

def line(s, x1,y1,x2,y2, color=LINE, w=1.5, dash=None):
    cn = s.shapes.add_connector(2, Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    cn.line.color.rgb = color; cn.line.width = Pt(w)
    cn.shadow.inherit = False
    if dash:
        ln = cn.line._get_or_add_ln()
        d = ln.makeelement(qn('a:prstDash'), {'val':dash}); ln.append(d)
    return cn

def arrow_conn(s, x1,y1,x2,y2, color=BLUE, w=2.0):
    cn = s.shapes.add_connector(2, Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    cn.line.color.rgb = color; cn.line.width = Pt(w); cn.shadow.inherit = False
    ln = cn.line._get_or_add_ln()
    tail = ln.makeelement(qn('a:tailEnd'), {'type':'triangle','w':'med','len':'med'})
    ln.append(tail)
    return cn

def _set_run(run, t, size, color, bold, font, italic=False):
    run.text = t
    f = run.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = font

def text(s, x,y,w,h, paras, anchor=MSO_ANCHOR.TOP, wrap=True):
    """paras: list of paragraphs. each = dict(runs=[(t,size,color,bold,font[,italic])...],
       align, sa(space_after pt), sb, line)."""
    tb = s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    tf.vertical_anchor = anchor
    for i,p in enumerate(paras):
        para = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        para.alignment = p.get('align', PP_ALIGN.LEFT)
        if 'sa' in p: para.space_after = Pt(p['sa'])
        if 'sb' in p: para.space_before = Pt(p['sb'])
        if 'line' in p: para.line_spacing = p['line']
        for r in p['runs']:
            t   = r[0]
            sz  = r[1] if len(r)>1 and r[1] else 16
            col = r[2] if len(r)>2 and r[2] else INK
            bd  = r[3] if len(r)>3 else False
            fn  = r[4] if len(r)>4 and r[4] else BODY
            it  = r[5] if len(r)>5 else False
            _set_run(para.add_run(), t, sz, col, bd, fn, it)
    return tb

def one(s,x,y,w,h, t, size=16, color=INK, bold=False, font=BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    return text(s,x,y,w,h,[{'runs':[(t,size,color,bold,font,italic)],'align':align}], anchor=anchor)

# section header (number chip + title)
def header(s, num, title):
    rect(s, 0.6,0.5,0.46,0.46, fill=BLUE, line=None, rad=0.25)
    one(s, 0.6,0.5,0.46,0.46, str(num), 18, WHITE, True, HEAD, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    one(s, 1.2,0.46,11.4,0.6, title, 30, NAVY, True, HEAD, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)

def subtitle(s, t):
    one(s, 1.2,1.12,11.5,0.4, t, 14.5, MUTED, False, BODY)

def footer(s, n):
    one(s, 0.6,7.06,7,0.3, "Dashboard EWS - PT Smart Home Inovasi", 10, RGBColor(0x9A,0xA7,0xB8), False, BODY)
    one(s, 9.4,7.06,3.33,0.3, "Dian Putri Iswandi - Sidang TA 2026", 10, RGBColor(0x9A,0xA7,0xB8), False, BODY, PP_ALIGN.RIGHT)
    one(s, 12.5,0.28,0.7,0.3, f"{n:02d}/15", 10, RGBColor(0xA9,0xB6,0xC8), False, MONO, PP_ALIGN.RIGHT)

# ---------- animation / transition ----------
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

def add_transition(s, kind="fade", spd="med"):
    from lxml import etree
    sld = s._element
    for el in sld.findall(qn('p:transition')): sld.remove(el)
    xml = f'<p:transition xmlns:p="{P_NS}" spd="{spd}"><p:{kind}/></p:transition>'
    node = etree.fromstring(xml)
    # place after clrMapOvr (before any timing)
    timing = sld.find(qn('p:timing'))
    if timing is not None: timing.addprevious(node)
    else: sld.append(node)

def add_entrance(s, spids, step=180, dur=420, first_on_click=False):
    """Auto-cascade fade entrance for the given shape ids (after previous)."""
    if not spids: return
    from lxml import etree
    cid = [3]
    def nxt():
        v = cid[0]; cid[0]+=1; return v
    pars = []
    for idx,sp in enumerate(spids):
        node = "clickEffect" if (idx==0 and first_on_click) else "afterEffect"
        outer_delay = "indefinite" if (idx==0 and first_on_click) else "0"
        eff_delay = "0" if idx==0 else str(step)
        a,b,c,d,e = nxt(),nxt(),nxt(),nxt(),nxt()
        pars.append(f'''<p:par><p:cTn id="{a}" fill="hold"><p:stCondLst><p:cond delay="{outer_delay}"/></p:stCondLst><p:childTnLst>
<p:par><p:cTn id="{b}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>
<p:par><p:cTn id="{c}" presetID="10" presetClass="entr" presetSubtype="0" fill="hold" grpId="0" nodeType="{node}"><p:stCondLst><p:cond delay="{eff_delay}"/></p:stCondLst><p:childTnLst>
<p:set><p:cBhvr><p:cTn id="{d}" dur="1" fill="hold"/><p:tgtEl><p:spTgt spid="{sp}"/></p:tgtEl><p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst></p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>
<p:animEffect transition="in" filter="fade"><p:cBhvr><p:cTn id="{e}" dur="{dur}"/><p:tgtEl><p:spTgt spid="{sp}"/></p:tgtEl></p:cBhvr></p:animEffect>
</p:childTnLst></p:cTn></p:par>
</p:childTnLst></p:cTn></p:par>
</p:childTnLst></p:cTn></p:par>''')
    timing_xml = f'''<p:timing xmlns:p="{P_NS}"><p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>
<p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>
{''.join(pars)}
</p:childTnLst></p:cTn>
<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>'''
    sld = s._element
    for el in sld.findall(qn('p:timing')): sld.remove(el)
    sld.append(etree.fromstring(timing_xml))

def pill(s, x,y,w,h, t, fill, txtcol, line=None):
    rect(s, x,y,w,h, fill=fill, line=line, rad=0.5)
    one(s, x,y,w,h, t, 12.5, txtcol, True, BODY, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

# ============================================================ SLIDE 1 COVER
s = slide()
rect(s, 0,0,0.32,7.5, fill=NAVY, line=None, rounded=False)  # left accent bar
rect(s, 0.9,0.85,0.62,0.62, fill=BLUE, line=None, rad=0.22)
one(s, 0.9,0.85,0.62,0.62, "SHI", 16, WHITE, True, HEAD, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
text(s, 1.7,0.86,9,0.7,[
  {'runs':[("UNIVERSITAS TEKNOLOGI YOGYAKARTA",12,MUTED,True,BODY)],'sa':1},
  {'runs':[("Prodi Sistem Informasi - Program Sarjana - 2026",11.5,MUTED,False,BODY)]},
])
one(s, 0.92,1.95,8,0.4, "SIDANG TUGAS AKHIR", 13, GOLD, True, BODY)
title_box = text(s, 0.9,2.4,11.6,1.8,[
  {'runs':[("Pengembangan Fitur ",30,NAVY,True,HEAD),("Dashboard",30,BLUE,True,HEAD),
           (" pada Sistem Manajemen Proyek ",30,NAVY,True,HEAD)],'line':1.05},
  {'runs':[("Berdasarkan ",30,NAVY,True,HEAD),("Data Laporan Harian",30,GOLD,True,HEAD)],'line':1.05,'sb':2},
])
one(s, 0.92,4.5,11.4,0.5, "Early Warning System (EWS) berbasis Schedule Performance Index  -  Studi Kasus: PT Smart Home Inovasi, Yogyakarta", 14.5, MUTED, False, BODY)
# RAG motif dots
for i,c in enumerate([GREEN,AMBER,RED]):
    dot(s, 0.92+i*0.34, 5.25, 0.2, c)
# meta
meta = [("PENYUSUN","Dian Putri Iswandi","NPM 5220311118"),
        ("DOSEN PEMBIMBING","Adityo Permana Wibowo, S.Kom., M.Cs.",""),
        ("STUDI KASUS","PT Smart Home Inovasi - IoT","")]
mx = 0.92
for k,v,subv in meta:
    text(s, mx,5.95,3.7,1.0,[
      {'runs':[(k,10.5,MUTED,True,BODY)],'sa':2},
      {'runs':[(v,14,INK,True,BODY)],'sa':1},
      {'runs':[(subv,11,MUTED,False,MONO)]},
    ])
    mx += 3.9
one(s, 12.5,0.28,0.7,0.3, "01/15", 10, RGBColor(0xA9,0xB6,0xC8), False, MONO, PP_ALIGN.RIGHT)
add_transition(s)
add_entrance(s, [title_box.shape_id])

# ============================================================ SLIDE 2 PERUSAHAAN
s = slide(); header(s,1,"Deskripsi Perusahaan")
c1 = rect(s, 0.6,1.35,5.55,4.7, fill=CARD)
rect(s, 0.95,1.7,0.42,0.42, fill=BLUE, line=None, rad=0.2)
one(s, 0.95,1.7,0.42,0.42, "SHI", 11, WHITE, True, HEAD, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
text(s, 1.5,1.66,4.4,0.7,[
  {'runs':[("PT Smart Home Inovasi",16,INK,True,HEAD)],'sa':1},
  {'runs':[("Sleman, Daerah Istimewa Yogyakarta",12,MUTED,False,BODY)]},
])
text(s, 0.95,2.6,4.85,1.4,[
  {'runs':[("Perusahaan teknologi ",14,INK,False,BODY),("Smart Home & Internet of Things (IoT)",14,INK,True,BODY),
           (": instalasi perangkat pintar, integrasi sistem, smart lighting & security untuk hunian dan bisnis.",14,INK,False,BODY)],'line':1.25},
])
chips2 = [("Instalasi IoT",CARD2),("Integrasi Sistem",CARD2),("Smart Lighting",CARD2),("Pengujian & Serah Terima",CARD2)]
cx,cy = 0.95,4.25
for t,f in chips2:
    w = 0.34+len(t)*0.082
    if cx+w>5.9: cx=0.95; cy+=0.55
    pill(s,cx,cy,w,0.4,t,CARD2,BLUED,line=LINE); cx+=w+0.18
# right: process chain
c2 = rect(s, 6.45,1.35,6.28,4.7, fill=WHITE)
one(s, 6.75,1.6,5.6,0.4, "Alur Proses Bisnis (per proyek)", 13.5, NAVY, True, BODY)
steps2 = ["Pesanan","Survei Lokasi","Desain","Instalasi","Integrasi","Pengujian","Serah Terima"]
bx,by,bw,bh = 6.72,2.25,0.58,0.74
anim2 = []
for i,st in enumerate(steps2):
    r = rect(s, bx,by,bw,bh, fill=CARD, line=LINE, rad=0.14)
    one(s, bx,by,bw,bh, str(i+1), 15, BLUE, True, HEAD, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    one(s, bx-0.14,by+bh+0.03,bw+0.28,0.45, st, 9, INK, False, BODY, PP_ALIGN.CENTER)
    anim2.append(r.shape_id)
    bx += bw
    if i<len(steps2)-1:
        chev(s, bx+0.005,by+0.24,0.13,0.26, MUTED); bx += 0.18
one(s, 6.75,3.95,5.7,1.2,
    "Setiap proyek dijalankan teknisi di lokasi klien dan dipantau manajer - inilah data operasional yang dikelola dashboard.",
    13, MUTED, False, BODY)
footer(s,2); add_transition(s); add_entrance(s,[c1.shape_id,c2.shape_id]+anim2, step=120)

# ============================================================ SLIDE 3 ANALISA KEBUTUHAN
s = slide(); header(s,2,"Analisa Kebutuhan"); subtitle(s,"Identifikasi masalah di lapangan & perumusan aliran data sistem usulan.")
one(s, 0.6,1.7,6,0.35,"IDENTIFIKASI MASALAH DI LAPANGAN",12.5,NAVY,True,BODY)
probs = [("1. Pelaporan tidak langsung","Teknisi lapor via WhatsApp, manajer input manual ke spreadsheet - teknisi tak punya akses sistem."),
         ("2. Tidak ada dashboard ringkas","Tak ada visualisasi kondisi seluruh proyek dalam satu layar."),
         ("3. Tanpa indikator risiko terukur","Deteksi keterlambatan reaktif, bukan preventif - proyek kritis terabaikan.")]
anim3 = []
py = 2.15
for t,d in probs:
    r = rect(s, 0.6,py,6.1,1.15, fill=CARD)
    rect(s, 0.6,py,0.09,1.15, fill=RED, line=None, rounded=False)
    text(s, 0.85,py+0.16,5.7,0.9,[
      {'runs':[(t,14.5,INK,True,BODY)],'sa':2},
      {'runs':[(d,12,MUTED,False,BODY)],'line':1.15},
    ])
    anim3.append(r.shape_id); py += 1.32
# right data flow
fc = rect(s, 7.05,1.7,5.68,4.55, fill=WHITE)
one(s, 7.35,1.95,5,0.35,"PERUMUSAN ALIRAN DATA",12.5,NAVY,True,BODY)
b1 = rect(s, 8.4,2.4,3,0.85, fill=CARD, line=LINE)
text(s,8.4,2.5,3,0.7,[{'runs':[("Teknisi",13,INK,True,BODY)],'align':PP_ALIGN.CENTER,'sa':1},{'runs':[("laporan harian",10.5,MUTED,False,BODY)],'align':PP_ALIGN.CENTER}])
b2 = rect(s, 8.4,3.75,3,0.92, fill=CARD2, line=BLUE)
text(s,8.4,3.86,3,0.75,[{'runs':[("Sistem / Modul EWS",13,BLUED,True,BODY)],'align':PP_ALIGN.CENTER,'sa':1},{'runs':[("hitung EV - PV - SPI",10.5,MUTED,False,BODY)],'align':PP_ALIGN.CENTER}])
b3 = rect(s, 8.4,5.2,3,0.85, fill=CARD, line=LINE)
text(s,8.4,5.3,3,0.7,[{'runs':[("Manajer",13,INK,True,BODY)],'align':PP_ALIGN.CENTER,'sa':1},{'runs':[("RAG + peringatan",10.5,MUTED,False,BODY)],'align':PP_ALIGN.CENTER}])
arrow_conn(s, 9.9,3.25,9.9,3.75, BLUE, 2)
arrow_conn(s, 9.9,4.67,9.9,5.2, BLUE, 2)
pill(s,7.35,5.55,1.0,0.0,"",None,None) if False else None
footer(s,3); add_transition(s); add_entrance(s, anim3+[fc.shape_id,b1.shape_id,b2.shape_id,b3.shape_id], step=130)

# ============================================================ SLIDE 4 SEBELUM
s = slide()
rect(s, 0.6,0.5,0.62,0.46, fill=RED, line=None, rad=0.18)
one(s, 0.6,0.5,0.62,0.46,"3a",15,WHITE,True,HEAD,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
text(s,1.35,0.46,11.4,0.55,[{'runs':[("Proses Bisnis ",28,NAVY,True,HEAD),("Sebelum",28,RED,True,HEAD),(" - Pelaporan Manual",28,NAVY,True,HEAD)],'align':PP_ALIGN.LEFT}],anchor=MSO_ANCHOR.MIDDLE)
subtitle(s,"Alur berjalan yang memicu keterlambatan deteksi masalah di PT SHI.")
steps4 = [("Teknisi lapor","via WhatsApp",CARD,INK),("Manajer baca","chat manual",CARD,INK),
          ("Input manual","ke spreadsheet",CARD,INK),("Hitung SPI","manual",CARD,INK),
          ("Deteksi","terlambat",RGBColor(0xFB,0xEA,0xE8),RED)]
bx,by,bw,bh = 0.6,2.0,2.18,1.35
anim4=[]
for i,(a,b,fl,tc) in enumerate(steps4):
    r = rect(s, bx,by,bw,bh, fill=fl, line=(RED if tc==RED else LINE))
    text(s,bx,by+0.42,bw,0.7,[{'runs':[(a,13.5,tc,True,BODY)],'align':PP_ALIGN.CENTER,'sa':1},{'runs':[(b,12,(RED if tc==RED else MUTED),False,BODY)],'align':PP_ALIGN.CENTER}])
    anim4.append(r.shape_id); bx+=bw
    if i<4: chev(s,bx+0.02,by+0.5,0.2,0.32,MUTED); bx+=0.27
co = rect(s, 0.6,3.75,12.13,1.0, fill=RGBColor(0xFB,0xEA,0xE8), line=RGBColor(0xF0,0xC4,0xBE))
text(s,0.9,3.95,11.6,0.7,[{'runs':[("Akibat:  ",14,RED,True,BODY),("rantai manual & berlapis membuat masalah baru ketahuan setelah terjadi - penanganan ",14,INK,False,BODY),("reaktif, bukan preventif",14,INK,True,BODY),(" -> keterlambatan proyek & menurunkan kepuasan pelanggan.",14,INK,False,BODY)],'line':1.2}])
cx=0.6
for t in ["Manual","Rawan human error","Tidak real-time","Teknisi tak terlibat"]:
    w=0.34+len(t)*0.085; pill(s,cx,5.05,w,0.42,t,RGBColor(0xFB,0xEA,0xE8),RED,line=RGBColor(0xF0,0xC4,0xBE)); cx+=w+0.2
footer(s,4); add_transition(s); add_entrance(s, anim4+[co.shape_id], step=140)

# ============================================================ SLIDE 5 SESUDAH
s = slide()
rect(s, 0.6,0.5,0.62,0.46, fill=GREEN, line=None, rad=0.18)
one(s, 0.6,0.5,0.62,0.46,"3b",15,WHITE,True,HEAD,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
text(s,1.35,0.46,11.4,0.55,[{'runs':[("Proses Bisnis ",28,NAVY,True,HEAD),("Sesudah",28,GREEN,True,HEAD),(" - Dashboard Terpusat",28,NAVY,True,HEAD)]}],anchor=MSO_ANCHOR.MIDDLE)
subtitle(s,"Dua peran, satu sumber kebenaran. Data mengalir otomatis - tanpa rantai manual.")
# teknisi card
tk = rect(s, 0.6,1.75,4.3,3.6, fill=CARD)
rect(s, 0.9,2.05,0.5,0.5, fill=BLUE, line=None, rad=0.2)
one(s,0.9,2.05,0.5,0.5,"T",18,WHITE,True,HEAD,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
text(s,0.9,2.7,3.7,0.6,[{'runs':[("Teknisi",17,INK,True,HEAD)],'sa':1},{'runs':[("Penginput Laporan Harian ",12,MUTED,False,BODY),("Mandiri",12,BLUE,True,BODY)]}])
for j,li in enumerate(["Input progres & kendala langsung dari lapangan","Kelola tugas di papan Kanban + unggah bukti","Ajukan eskalasi kendala lapangan"]):
    dot(s,0.95,3.62+j*0.5,0.1,BLUE); one(s,1.2,3.5+j*0.5,3.45,0.5,li,12,INK,False,BODY)
# manajer card
mn = rect(s, 8.43,1.75,4.3,3.6, fill=CARD)
rect(s, 8.73,2.05,0.5,0.5, fill=GOLD, line=None, rad=0.2)
one(s,8.73,2.05,0.5,0.5,"M",18,WHITE,True,HEAD,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
text(s,8.73,2.7,3.7,0.6,[{'runs':[("Manajer",17,INK,True,HEAD)],'sa':1},{'runs':[("Pengawas ",12,MUTED,False,BODY),("Dashboard Sentral",12,GOLD,True,BODY)]}])
for j,li in enumerate(["Pantau seluruh proyek aktif dalam satu layar","Terima peringatan dini - proyek kritis teratas","Review bukti, setujui tugas, balas eskalasi"]):
    dot(s,8.78,3.62+j*0.5,0.1,GOLD); one(s,9.03,3.5+j*0.5,3.45,0.5,li,12,INK,False,BODY)
# center dashboard
cd = oval(s, 5.46,2.55,2.4,2.0, fill=CARD2, line=BLUE)
text(s,5.46,3.0,2.4,0.8,[{'runs':[("Dashboard",13,BLUED,True,BODY)],'align':PP_ALIGN.CENTER,'sa':1},{'runs':[("EWS Sentral",11,MUTED,False,BODY)],'align':PP_ALIGN.CENTER}])
for k,c in enumerate([GREEN,AMBER,RED]): dot(s,6.16+k*0.34,3.95,0.18,c)
arrow_conn(s,4.9,3.5,5.46,3.5,MUTED,1.8)
arrow_conn(s,7.86,3.5,8.43,3.5,MUTED,1.8)
co5 = rect(s, 0.6,5.55,12.13,0.85, fill=RGBColor(0xE7,0xF4,0xEC), line=RGBColor(0xBF,0xE3,0xCD))
text(s,0.9,5.72,11.6,0.6,[{'runs':[("Hasil:  ",14,GREEN,True,BODY),("teknisi melapor sendiri secara real-time, manajer memantau terpusat - deteksi masalah jadi ",14,INK,False,BODY),("preventif",14,INK,True,BODY),(".",14,INK,False,BODY)]}])
footer(s,5); add_transition(s); add_entrance(s,[tk.shape_id,cd.shape_id,mn.shape_id,co5.shape_id], step=150)

# ============================================================ SLIDE 6 ARSITEKTUR
s = slide(); header(s,4,"Pemodelan - Arsitektur Sistem"); subtitle(s,"Arsitektur berlapis: antarmuka -> logika komputasi EWS -> basis data.")
layers = [("Klien",BLUE,"Pengguna","Teknisi & Manajer (peramban / browser)",None),
          ("Frontend",BLUE,"Next.js","(SSR, React) + Tailwind CSS + Recharts - render dashboard & Kanban",BLUE),
          ("Logika",GOLD,"Modul Komputasi EWS","API + JWT Auth + RBAC + perhitungan SPI, PV/EV, klasifikasi RAG, criticality sort",GOLD),
          ("Data",GREEN,"PostgreSQL","7 tabel (user, klien, proyek, tugas, bukti, eskalasi, penugasan) + penyimpanan bukti (disk)",GREEN)]
ly=1.75; anim6=[]
for i,(tag,tagc,h1,h2,bord) in enumerate(layers):
    r = rect(s, 0.6,ly,12.13,0.96, fill=CARD, line=(bord or LINE), lw=(1.6 if bord else 1.0))
    pill(s,0.9,ly+0.28,1.5,0.4,tag,tagc,WHITE)
    text(s,2.65,ly+0.18,9.8,0.62,[{'runs':[(h1+"  ",15,INK,True,BODY),(h2,13,MUTED,False,BODY)],'line':1.1}],anchor=MSO_ANCHOR.MIDDLE)
    anim6.append(r.shape_id); ly+=0.96
    if i<3:
        line(s,6.66,ly,6.66,ly+0.22,MUTED,1.5); ly+=0.22
footer(s,6); add_transition(s); add_entrance(s,anim6, step=150)

# ============================================================ SLIDE 7 USECASE
s = slide(); header(s,4,"Use Case Diagram")
subtitle(s,"Dua aktor - Teknisi & Manajer. Fungsi inti wajib melalui <<include>> Login.")
# boundary
rect(s, 2.0,1.55,9.3,4.65, fill=None, line=LINE, rad=0.02)
one(s, 2.0,1.62,9.3,0.3,"Sistem Dashboard EWS",11,MUTED,False,MONO,PP_ALIGN.CENTER)
def actor(s,x,y,name,color):
    oval(s,x+0.13,y,0.26,0.26,fill=None,line=color,lw=2)        # head
    line(s,x+0.26,y+0.26,x+0.26,y+0.72,color,2)                 # body
    line(s,x+0.05,y+0.4,x+0.47,y+0.4,color,2)                   # arms
    line(s,x+0.26,y+0.72,x+0.1,y+1.0,color,2)                   # leg
    line(s,x+0.26,y+0.72,x+0.42,y+1.0,color,2)
    one(s,x-0.25,y+1.02,0.76,0.3,name,12,color,True,BODY,PP_ALIGN.CENTER)
actor(s,0.85,3.3,"Teknisi",BLUED)
actor(s,11.55,3.3,"Manajer",GOLD)
def uc(s,x,y,t,fill=CARD,line_=LINE,tc=INK,w=2.55):
    o = oval(s,x,y,w,0.62,fill=fill,line=line_,lw=1.3)
    one(s,x,y,w,0.62,t,11.5,tc,False,BODY,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
    return o,(x+w/2,y+0.31),x,x+w
# login center
lg,lgc,_,_ = uc(s,5.39,3.34,"Login",CARD2,BLUE,BLUED,2.0)
tk_uc = [("Tinjau dashboard performa",1.9),("Lihat detail tugas & proyek",2.55),("Mengisi daily report",4.25),("Mengajukan eskalasi",4.9)]
tkpts=[]
for t,y in tk_uc:
    o,c,xl,xr = uc(s,2.25,y,t); tkpts.append((c,xl))
mn_uc = [("Tinjau dashboard proyek",1.9),("Kelola data proyek",2.55),("Kelola penugasan teknisi",4.25),("Menindaklanjuti eskalasi",4.9)]
mnpts=[]
for t,y in mn_uc:
    o,c,xl,xr = uc(s,6.5,y,t); mnpts.append((c,xr))
# exclude
oe = oval(s,6.5,5.55,2.55,0.55, fill=RGBColor(0xFB,0xEA,0xE8), line=RED, lw=1.2)
one(s,6.5,5.55,2.55,0.55,"Hapus data pelanggan <<exclude>>",10,RED,False,BODY,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
# connectors actor->uc
for c,xl in tkpts: line(s,1.35,3.75,xl,c[1],LINE,1.0)
for c,xr in mnpts: line(s,11.85,3.75,xr,c[1],LINE,1.0)
# legend
lx=0.6
for t,fc,tc in [("<<include>> wajib Login",CARD2,BLUED),("<<extend>> riwayat / detail",CARD,MUTED),("<<exclude>> hapus pelanggan",RGBColor(0xFB,0xEA,0xE8),RED)]:
    w=0.3+len(t)*0.082; pill(s,lx,6.55,w,0.4,t,fc,tc,line=LINE); lx+=w+0.2
footer(s,7); add_transition(s)
add_entrance(s,[lg.shape_id], step=150)

# ============================================================ SLIDE 8 WIREFRAME MANAJER
s = slide(); header(s,4,"Wireframe - Dashboard Manajer (EWS)")
mock = rect(s, 0.6,1.4,12.13,4.85, fill=WHITE, line=LINE)
rect(s, 0.6,1.4,12.13,0.5, fill=NAVY, line=None, rounded=False)
one(s,0.85,1.4,5,0.5,"PT Smart Home Inovasi",12,WHITE,True,BODY,PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)
one(s,8,1.4,4.5,0.5,"Manajer: Budi   |   Keluar",11,RGBColor(0xCF,0xDA,0xEC),False,BODY,PP_ALIGN.RIGHT,MSO_ANCHOR.MIDDLE)
# sidebar
rect(s, 0.6,1.9,1.7,4.35, fill=RGBColor(0xF4,0xF7,0xFB), line=LINE, rounded=False)
one(s,0.78,2.05,1.4,0.3,"MENU",9,MUTED,True,BODY)
for j,(m,on) in enumerate([("Dashboard",True),("Proyek",False),("Klien",False),("Jadwal",False),("Laporan",False)]):
    if on: rect(s,0.7,2.35+j*0.42,1.5,0.36,fill=CARD2,line=None,rad=0.1)
    one(s,0.82,2.35+j*0.42,1.4,0.36,("# " if on else "")+m,11,(INK if on else MUTED),on,BODY,PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)
# stat cards
cards = [("TOTAL PROYEK","12",INK,LINE),("STATUS MERAH","3",RED,RED),("STATUS KUNING","4",AMBER,AMBER),("STATUS HIJAU","5",GREEN,GREEN)]
sx=2.5; anim8=[]
for k,(lbl,v,vc,bd) in enumerate(cards):
    r=rect(s,sx,2.1,2.42,1.0,fill=CARD,line=bd,lw=(1.4 if bd!=LINE else 1.0))
    one(s,sx+0.18,2.22,2.1,0.3,lbl,9.5,MUTED,True,BODY)
    one(s,sx+0.16,2.5,2.1,0.5,v,24,vc,True,HEAD)
    anim8.append(r.shape_id); sx+=2.5
one(s,2.5,3.25,8,0.3,"Daftar Proyek - diurutkan dari status paling kritis",10.5,MUTED,False,BODY)
# table
tbl_rows=[("STATUS","PROYEK","KLIEN","SPI","DEVIASI",None),
          ("","Instalasi Cluster A","PT Mitra Aksara","0.72","-28%",RED),
          ("","Renovasi Server B","PT Bina Daya","0.81","-19%",RED),
          ("","Pemasangan IoT Lt 3","Yayasan Cahaya","0.88","-12%",AMBER),
          ("","Instalasi Smart Office","PT Maju Bersama","0.97","-3%",GREEN)]
ty=3.6
for ri,(a,b,c,d,e,rc) in enumerate(tbl_rows):
    if ri==0:
        one(s,2.7,ty,1.0,0.3,a,9.5,MUTED,True,BODY); one(s,3.7,ty,3.0,0.3,b,9.5,MUTED,True,BODY)
        one(s,6.9,ty,2.4,0.3,c,9.5,MUTED,True,BODY); one(s,9.5,ty,1.2,0.3,d,9.5,MUTED,True,BODY); one(s,10.7,ty,1.6,0.3,e,9.5,MUTED,True,BODY)
        line(s,2.5,ty+0.32,12.4,ty+0.32,LINE,1.0); ty+=0.42; continue
    dot(s,2.78,ty+0.04,0.18,rc)
    one(s,3.7,ty,3.0,0.32,b,12,INK,False,BODY); one(s,6.9,ty,2.4,0.32,c,12,INK,False,BODY)
    one(s,9.5,ty,1.2,0.32,d,12,rc,True,MONO); one(s,10.7,ty,1.6,0.32,e,12,rc,False,BODY)
    line(s,2.5,ty+0.38,12.4,ty+0.38,RGBColor(0xEE,0xF2,0xF7),0.75); ty+=0.46
text(s,2.5,5.95,9.9,0.3,[{'runs':[("Eskalasi terbaru:  ",10.5,RED,True,BODY),("Andi (Teknisi) - Cluster A: kabel utama tidak kompatibel - 07/05 08:50",10.5,MUTED,False,BODY)]}])
footer(s,8); add_transition(s); add_entrance(s,[mock.shape_id]+anim8, step=110)

# ============================================================ SLIDE 9 WIREFRAME TEKNISI
s = slide(); header(s,4,"Wireframe - Dashboard Teknisi (Performa)")
kpis=[("SPI SAYA","0.94",AMBER,"target >= 0.95"),("TUGAS SELESAI","18",INK,"30 hari"),
      ("TUGAS BERJALAN","4",INK,"saat ini"),("OVERDUE","1",RED,"perlu eskalasi")]
kx=0.6; anim9=[]
for lbl,v,vc,sub in kpis:
    r=rect(s,kx,1.5,2.78,1.15,fill=CARD,line=(RED if vc==RED else LINE),lw=(1.3 if vc==RED else 1.0))
    one(s,kx+0.18,1.62,2.5,0.3,lbl,9.5,MUTED,True,BODY)
    one(s,kx+0.16,1.9,2.5,0.5,v,24,vc,True,HEAD)
    one(s,kx+0.18,2.36,2.5,0.25,sub,10,MUTED,False,MONO)
    anim9.append(r.shape_id); kx+=2.92
# donut composition (chart) bottom-left
dch = rect(s,0.6,2.9,5.6,3.25,fill=WHITE,line=LINE)
one(s,0.85,3.1,5,0.3,"Komposisi Status Tugas (30 hari)",12.5,NAVY,True,BODY)
cd2 = CategoryChartData(); cd2.categories=["Selesai","To Do","Working/Review"]; cd2.add_series("s",(18,4,2))
gf = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(0.8),Inches(3.45),Inches(2.5),Inches(2.5), cd2)
ch=gf.chart; ch.has_legend=False; ch.has_title=False
for i,c in enumerate([GREEN,BLUE,AMBER]):
    pt=ch.plots[0].series[0].points[i]; pt.format.fill.solid(); pt.format.fill.fore_color.rgb=c
for j,(t,c) in enumerate([("Selesai (18)",GREEN),("To Do (4)",BLUE),("Working / Review (2)",AMBER)]):
    dot(s,3.55,3.7+j*0.45,0.16,c); one(s,3.8,3.62+j*0.45,2.3,0.35,t,11.5,INK,False,BODY)
# kanban
kb = rect(s,6.4,2.9,6.33,3.25,fill=WHITE,line=LINE)
one(s,6.65,3.1,5,0.3,"Papan Tugas Saya (Kanban)",12.5,NAVY,True,BODY)
cols=[("TO DO",["Pasang sensor utama|07/05","Tarik kabel lantai 2|08/05"]),
      ("WORKING ON IT",["Konfigurasi gateway|+ bukti foto"]),
      ("DONE  [manajer]",["Uji koneksi awal|disetujui"])]
cxx=6.6
for ci,(h,cards_) in enumerate(cols):
    rect(s,cxx,3.5,1.95,2.45,fill=RGBColor(0xF6,0xF9,0xFC),line=LINE,rad=0.05)
    one(s,cxx+0.1,3.58,1.8,0.3,h,9,(MUTED if ci<2 else AMBER),True,BODY)
    yy=3.95
    for cc in cards_:
        ti,me=cc.split("|")
        rect(s,cxx+0.1,yy,1.75,0.6,fill=WHITE,line=LINE,rad=0.08)
        text(s,cxx+0.2,yy+0.07,1.6,0.5,[{'runs':[(ti,10,INK,False,BODY)],'sa':1},{'runs':[(me,9,MUTED,False,MONO)]}])
        yy+=0.7
    cxx+=2.1
text(s,6.65,6.25,6.0,0.4,[{'runs':[("Review Gate:  ",11.5,GOLD,True,BODY),("teknisi geser to_do <-> working_on_it; status ",11.5,INK,False,BODY),("'Done' hanya disetujui Manajer.",11.5,INK,True,BODY)]}])
footer(s,9); add_transition(s); add_entrance(s,anim9+[dch.shape_id,kb.shape_id], step=120)

# ============================================================ SLIDE 10 IMPLEMENTASI
s = slide(); header(s,5,"Implementasi (Coding)")
one(s,0.6,1.55,6,0.3,"TUMPUKAN TEKNOLOGI",12.5,NAVY,True,BODY)
techs=["Next.js","React","TypeScript","Tailwind CSS","Recharts","Node.js","PostgreSQL","JWT"]
cx=0.6;cy=1.95
for t in techs:
    w=0.36+len(t)*0.105
    if cx+w>6.4: cx=0.6;cy+=0.55
    rect(s,cx,cy,w,0.42,fill=CARD,line=LINE,rad=0.12); one(s,cx,cy,w,0.42,t,12,NAVY,True,MONO,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE); cx+=w+0.16
one(s,0.6,3.15,6,0.3,"MODUL TERBANGUN",12.5,NAVY,True,BODY)
mods=[("Autentikasi + RBAC",BLUE,"login JWT, pisah hak Teknisi/Manajer"),
      ("Kanban + Review Gate",BLUE,"ubah status tugas, unggah bukti, approve done"),
      ("Modul Komputasi EWS",GOLD,"SPI, PV/EV, klasifikasi RAG, urut kritis"),
      ("Dashboard Analitik",GREEN,"grafik Recharts + eskalasi")]
my=3.55; anim10=[]
for h,c,d in mods:
    r=rect(s,0.6,my,5.9,0.62,fill=CARD,line=LINE)
    rect(s,0.6,my,0.09,0.62,fill=c,line=None,rounded=False)
    text(s,0.85,my+0.07,5.5,0.5,[{'runs':[(h+"  ",13.5,INK,True,BODY),(d,11.5,MUTED,False,BODY)],'line':1.05}],anchor=MSO_ANCHOR.MIDDLE)
    anim10.append(r.shape_id); my+=0.72
# image
img = None
imgpath = os.path.join(os.path.dirname(__file__),"assets","charts-impl.png")
ic = rect(s,6.85,1.55,5.88,4.6,fill=WHITE,line=LINE)
if os.path.exists(imgpath):
    img = s.shapes.add_picture(imgpath, Inches(7.0),Inches(1.85),width=Inches(5.58))
one(s,7.0,5.55,5.6,0.5,"Implementasi nyata: dashboard analitik - beban kerja teknisi, distribusi & tenggat tugas (Recharts).",10.5,MUTED,False,BODY)
footer(s,10); add_transition(s); add_entrance(s,anim10+[ic.shape_id], step=120)

# ============================================================ SLIDE 11 ALGORITMA SPI
s = slide(); header(s,6,"Algoritma: Laporan Harian -> SPI -> Indikator RAG")
calc=[("1","Input - Laporan harian","status tugas berubah",INK),
      ("2","Earned Value (EV) - % tugas selesai","11 / 20 = 55%",GREEN),
      ("3","Planned Value (PV) - hari berjalan / durasi","54 / 90 = 60%",BLUE),
      ("4","SPI = EV / PV","0.55 / 0.60 = 0.92",GOLD)]
cy=1.55; anim11=[]
for n,k,v,vc in calc:
    r=rect(s,0.6,cy,6.4,0.78,fill=CARD,line=(AMBER if n=="4" else LINE),lw=(1.5 if n=="4" else 1.0))
    rect(s,0.78,cy+0.19,0.42,0.42,fill=(AMBER if n=="4" else BLUE),line=None,rad=0.22)
    one(s,0.78,cy+0.19,0.42,0.42,n,15,WHITE,True,HEAD,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
    one(s,1.4,cy+0.12,3.6,0.55,k,12,MUTED,False,BODY,PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)
    one(s,5.0,cy+0.12,1.85,0.55,v,15,vc,True,MONO,PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)
    anim11.append(r.shape_id); cy+=0.92
# formula box
fb=rect(s,0.6,5.35,6.4,1.05,fill=RGBColor(0xF4,0xF7,0xFB),line=LINE)
text(s,0.8,5.5,6.0,0.8,[
  {'runs':[("SPI = (tugas selesai/total) / (hari berjalan/durasi)",15,INK,True,MONO)],'align':PP_ALIGN.CENTER,'sa':3},
  {'runs':[("= (11/20) / (54/90) = ",15,INK,False,MONO),("0.92",16,GOLD,True,MONO)],'align':PP_ALIGN.CENTER}])
# RAG mapping
one(s,7.35,1.55,5,0.3,"5. PEMETAAN RAG (OTOMATIS)",12.5,NAVY,True,BODY)
rags=[("SPI >= 0.95","Hijau - on track",GREEN,False),
      ("0.85 <= SPI < 0.95","Kuning - peringatan",AMBER,True),
      ("SPI < 0.85","Merah - kritis",RED,False)]
ry=2.0; anim11b=[]
for cond,mean,c,hit in rags:
    r=rect(s,7.35,ry,5.38,0.72,fill=(RGBColor(0xFC,0xF4,0xE2) if hit else CARD),line=(AMBER if hit else LINE),lw=(1.6 if hit else 1.0))
    dot(s,7.6,ry+0.26,0.2,c)
    one(s,8.0,ry,2.3,0.72,cond,13,INK,False,MONO,PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)
    one(s,10.2,ry,2.4,0.72,mean,12,(AMBER if hit else MUTED),hit,BODY,PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)
    anim11b.append(r.shape_id); ry+=0.85
co11=rect(s,7.35,4.7,5.38,1.7,fill=RGBColor(0xFC,0xF4,0xE2),line=RGBColor(0xEB,0xD7,0xA8))
text(s,7.6,4.9,4.9,1.35,[{'runs':[("SPI ",13.5,INK,False,BODY),("0.92",13.5,GOLD,True,BODY),(" jatuh di rentang kuning -> sistem otomatis menandai proyek ",13.5,INK,False,BODY),("PERINGATAN",13.5,AMBER,True,BODY),(" & menaikkan urutannya di dashboard. Tanpa input persen manual.",13.5,INK,False,BODY)],'line':1.25}])
footer(s,11); add_transition(s); add_entrance(s,anim11+[fb.shape_id]+anim11b+[co11.shape_id], step=160)

# ============================================================ SLIDE 12 PENGUJIAN
s = slide(); header(s,7,"Pengujian (Testing)")
subtitle(s,"Metode Black Box - fokus kesesuaian input/output pada skenario interaksi Manajer & Teknisi.")
cats=[("3","TC-AUTH","Autentikasi & RBAC",BLUE),("5","TC-MN","Skenario Manajer",GOLD),
      ("6","TC-TK","Skenario Teknisi",GREEN),("4","TC-SYS","Komputasi Sistem",RGBColor(0x8A,0x5C,0xD6))]
cx=0.6; anim12=[]
for v,t,d,c in cats:
    r=rect(s,cx,1.75,2.93,1.2,fill=CARD,line=LINE)
    one(s,cx,1.9,2.93,0.5,v,30,c,True,HEAD,PP_ALIGN.CENTER)
    text(s,cx,2.42,2.93,0.5,[{'runs':[(t,11.5,INK,True,BODY)],'align':PP_ALIGN.CENTER,'sa':1},{'runs':[(d,10.5,MUTED,False,BODY)],'align':PP_ALIGN.CENTER}])
    anim12.append(r.shape_id); cx+=3.05
rows=[("TC-AUTH-03","Teknisi diblokir akses dashboard manajer via URL (RBAC)",False),
      ("TC-MN-03","Review gate: manajer tinjau bukti -> status Done",False),
      ("TC-SYS-01","Hitung EV otomatis dari bobot penugasan disetujui",False),
      ("TC-SYS-02","Hitung SPI (EV vs PV) -> render warna RAG",True),
      ("TC-SYS-03","Agregasi + urutkan proyek kritis ke atas",False)]
tcard=rect(s,0.6,3.25,12.13,2.35,fill=WHITE,line=LINE)
ty=3.45
one(s,0.85,ty,1.8,0.3,"ID",10,MUTED,True,BODY); one(s,2.8,ty,7.5,0.3,"SKENARIO KUNCI",10,MUTED,True,BODY); one(s,11.0,ty,1.5,0.3,"HASIL",10,MUTED,True,BODY,PP_ALIGN.CENTER)
line(s,0.85,ty+0.33,12.5,ty+0.33,LINE,1.0); ty+=0.45
for idn,desc,hl in rows:
    one(s,0.85,ty,1.9,0.32,idn,11,(GOLD if hl else MUTED),hl,MONO)
    one(s,2.8,ty,7.6,0.32,desc,12,INK,hl,BODY)
    pill(s,11.0,ty-0.03,1.5,0.36,"Valid",RGBColor(0xE7,0xF4,0xEC),GREEN,line=RGBColor(0xBF,0xE3,0xCD))
    line(s,0.85,ty+0.4,12.5,ty+0.4,RGBColor(0xEE,0xF2,0xF7),0.75); ty+=0.5
co12=rect(s,0.6,5.75,12.13,0.75,fill=RGBColor(0xE7,0xF4,0xEC),line=RGBColor(0xBF,0xE3,0xCD))
text(s,0.85,5.92,11.6,0.5,[{'runs':[("Seluruh skenario inti komputasi (TC-SYS-01..04) ",13.5,INK,False,BODY),("Valid",13.5,GREEN,True,BODY),(" - mesin EWS terbukti menghitung & mewarnai dengan benar.",13.5,INK,False,BODY)]}])
footer(s,12); add_transition(s); add_entrance(s,anim12+[tcard.shape_id,co12.shape_id], step=120)

# ============================================================ SLIDE 13 METRIK
s = slide(); header(s,8,"Metrik Keberhasilan")
# doughnut 83.33
cd3=CategoryChartData(); cd3.categories=["Valid","Tidak Valid"]; cd3.add_series("s",(15,3))
gf3=s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(0.8),Inches(1.5),Inches(3.0),Inches(3.0), cd3)
ch3=gf3.chart; ch3.has_legend=False; ch3.has_title=False
ch3.plots[0].series[0].points[0].format.fill.solid(); ch3.plots[0].series[0].points[0].format.fill.fore_color.rgb=GREEN
ch3.plots[0].series[0].points[1].format.fill.solid(); ch3.plots[0].series[0].points[1].format.fill.fore_color.rgb=RGBColor(0xE7,0xDF,0xD0)
one(s,0.8,2.6,3.0,0.55,"83,33%",23,GREEN,True,HEAD,PP_ALIGN.CENTER)
one(s,0.8,3.15,3.0,0.3,"tingkat valid",11,MUTED,False,BODY,PP_ALIGN.CENTER)
stat=[("18","Total Skenario",INK,LINE),("15","Valid",GREEN,GREEN),("3","Tidak Valid",RED,RED)]
sx=4.45; anim13=[]
for v,l,c,bd in stat:
    r=rect(s,sx,1.75,2.6,1.3,fill=CARD,line=bd,lw=(1.4 if bd!=LINE else 1.0))
    one(s,sx,1.95,2.6,0.6,v,32,c,True,HEAD,PP_ALIGN.CENTER)
    one(s,sx,2.62,2.6,0.3,l,12,MUTED,False,BODY,PP_ALIGN.CENTER)
    anim13.append(r.shape_id); sx+=2.75
co13=rect(s,4.45,3.25,8.28,1.0,fill=RGBColor(0xE7,0xF4,0xEC),line=RGBColor(0xBF,0xE3,0xCD))
text(s,4.7,3.42,7.8,0.75,[{'runs':[("Fitur komputasi EWS berfungsi akurat.  ",13.5,GREEN,True,BODY),("Ketiga temuan minor bukan pada inti perhitungan EWS, melainkan penyempurnaan UX/keamanan - menjadi bahan pemeliharaan korektif.",13,INK,False,BODY)],'line':1.2}])
one(s,0.6,4.55,8,0.3,"3 TEMUAN MINOR (-> PERBAIKAN)",12.5,NAVY,True,BODY)
finds=[("TC-MN-04","Indikator EWS belum real-time - butuh refresh manual."),
       ("TC-TK-03","Validasi unggah berkas (.exe) masih lolos."),
       ("TC-TK-06","Klik ganda \"Kirim\" menyebabkan laporan duplikat.")]
fx=0.6
for idn,d in finds:
    rect(s,fx,4.95,3.95,1.15,fill=CARD,line=LINE)
    rect(s,fx,4.95,0.09,1.15,fill=AMBER,line=None,rounded=False)
    text(s,fx+0.25,5.1,3.55,0.9,[{'runs':[(idn,11,GOLD,True,MONO)],'sa':2},{'runs':[(d,12,INK,False,BODY)],'line':1.15}])
    fx+=4.07
footer(s,13); add_transition(s); add_entrance(s,anim13+[co13.shape_id], step=130)

# ============================================================ SLIDE 14 PEMELIHARAAN
s = slide(); header(s,9,"Pemeliharaan (Maintenance)")
subtitle(s,"Empat jenis pemeliharaan perangkat lunak untuk menjaga sistem tetap andal & relevan.")
maint=[("Korektif",RED,"Perbaiki 3 temuan uji: refresh real-time, validasi unggah, anti double-submit."),
       ("Adaptif",BLUE,"Selaras dengan perubahan versi Next.js, PostgreSQL, peramban & sistem operasi."),
       ("Perfektif",GOLD,"Tingkatkan fitur: notifikasi otomatis, aplikasi mobile teknisi, dashboard biaya."),
       ("Preventif",GREEN,"Backup berkala, pemantauan kinerja, audit keamanan & RBAC.")]
mx=0.6; anim14=[]
for h,c,d in maint:
    r=rect(s,mx,1.85,2.93,2.3,fill=CARD,line=LINE)
    rect(s,mx,1.85,2.93,0.1,fill=c,line=None,rounded=False)
    text(s,mx+0.22,2.15,2.5,1.8,[{'runs':[(h,16,INK,True,HEAD)],'sa':6},{'runs':[(d,12.5,MUTED,False,BODY)],'line':1.25}])
    anim14.append(r.shape_id); mx+=3.06
sc=rect(s,0.6,4.5,12.13,1.4,fill=WHITE,line=LINE)
one(s,0.85,4.7,8,0.3,"SARAN PENGEMBANGAN (BAB VI)",12.5,NAVY,True,BODY)
sx=0.85
for t,c in [("Aplikasi mobile untuk teknisi",BLUE),("Notifikasi otomatis (push)",GOLD),("Integrasi aspek biaya proyek",GREEN)]:
    w=0.4+len(t)*0.095; rect(s,sx,5.15,w,0.5,fill=CARD,line=c,lw=1.3,rad=0.5); one(s,sx,5.15,w,0.5,t,12,c,True,BODY,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE); sx+=w+0.25
footer(s,14); add_transition(s); add_entrance(s,anim14+[sc.shape_id], step=130)

# ============================================================ SLIDE 15 PENUTUP
s = slide()
rect(s, 0,0,0.32,7.5, fill=NAVY, line=None, rounded=False)
one(s,0.9,0.7,5,0.4,"PENUTUP",13,GOLD,True,BODY)
tb15=text(s,0.9,1.2,11.6,1.6,[
  {'runs':[("Dashboard berhasil mengubah ",29,NAVY,True,HEAD),("data laporan harian",29,GOLD,True,HEAD),(" menjadi",29,NAVY,True,HEAD)],'line':1.08},
  {'runs':[("indikator visual kesehatan proyek",29,BLUE,True,HEAD),(" secara otomatis.",29,NAVY,True,HEAD)],'line':1.08,'sb':2}])
ca=rect(s,0.9,3.2,5.85,1.5,fill=RGBColor(0xE7,0xF4,0xEC),line=RGBColor(0xBF,0xE3,0xCD))
text(s,1.15,3.4,5.4,1.2,[{'runs':[("Tujuan tercapai.  ",13.5,GREEN,True,BODY),("Deteksi keterlambatan kini preventif - SPI -> RAG -> proyek kritis tampil teratas. Pengujian 83,33% valid; inti EWS akurat.",13,INK,False,BODY)],'line':1.25}])
cb=rect(s,7.0,3.2,5.73,1.5,fill=CARD,line=LINE)
text(s,7.25,3.4,5.3,1.2,[{'runs':[("Kontribusi:  ",13.5,NAVY,True,BODY),("akses langsung teknisi, dashboard EWS terpusat, komputasi SPI otomatis, peringatan dini berbasis data.",13,MUTED,False,BODY)],'line':1.25}])
one(s,0.9,5.5,8,0.9,"Terima Kasih",40,NAVY,True,HEAD)
one(s,0.95,6.45,5,0.3,"Sesi Tanya Jawab",13,MUTED,False,BODY)
text(s,8.5,6.05,4.2,0.8,[{'runs':[("Dian Putri Iswandi",14,INK,True,BODY)],'align':PP_ALIGN.RIGHT,'sa':1},{'runs':[("5220311118 - UTY 2026",11,MUTED,False,MONO)],'align':PP_ALIGN.RIGHT}])
for i,c in enumerate([GREEN,AMBER,RED]): dot(s,0.92,5.0+0,0.0,c) if False else dot(s,0.92+i*0.34,5.05,0.2,c)
one(s, 12.5,0.28,0.7,0.3, "15/15", 10, RGBColor(0xA9,0xB6,0xC8), False, MONO, PP_ALIGN.RIGHT)
add_transition(s); add_entrance(s,[tb15.shape_id,ca.shape_id,cb.shape_id], step=150)

out = os.path.join(os.path.dirname(__file__),"presentasi-sidang.pptx")
prs.save(out)
print("SAVED", out, "slides:", len(prs.slides._sldIdLst))
